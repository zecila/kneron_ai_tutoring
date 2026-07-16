import json
from pathlib import Path
import numpy as np
import pdfplumber

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from lxml import etree
from collections import Counter
from statistics import mean

import pytesseract
import cv2
from docx import Document
from docx.text.paragraph import Paragraph
from docx.table import Table
import os
import re

from pipeline.ocr_pipeline import extract_image_elements, extract_image_elements_pptx
from pdf2image import convert_from_path

def clean_text(text):
  lines = text.splitlines()
  cleaned = []

  for line in lines:
    line = line.strip()

    if not line:
      continue
    if len(line) < 2:
      continue

    cleaned.append(line)

  # DO NOT TOUCH!! OTHERWISE TABLE CONTENTS NOT EXTRACTED
  return re.sub(r'\s+', " ", text).strip()

def get_bbox(shape):
  return {
    "x": int(shape.left),
    "y": int(shape.top),
    "width": int(shape.width),
    "height": int(shape.height)
  }

#------------------------------------------------
# HANDLING EQUATIONS AND SPECIAL CHARACTERS
#------------------------------------------------

OMML_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
XSLT_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "OMML2MML.XSL")
MATH_SYMBOLS = set("∑∏∫∂∇√∞≠≤≥±×÷∈∉⊂⊃∩∪⇒⇔→←αβγδεθλμπσφψω")
SYMBOL_FONT_MATH = set("\uf028\uf029\uf02b\uf02d\uf03d\uf040\uf044\uf045\uf046\uf047\uf048\uf049\uf04a\uf04b\uf04c\uf04d\uf04e\uf050\uf051\uf052\uf053\uf054\uf055\uf056\uf057\uf058\uf059\uf05a\uf0a3\uf0b1\uf0b4\uf0b8\uf0d0\uf0d1\uf0e5\uf0e6\uf0f2\uf0f3\uf0f4\uf0f5")
_xslt_transform = None

SYMBOL_FONT_MAP = {
    "\uf053": "\\Sigma",
    "\uf050": "\\pi",
    "\uf044": "\\Delta",
    "\uf06d": "\\mu",
    "\uf062": "\\beta",
    "\uf061": "\\alpha",
    "\uf067": "\\gamma",
    "\uf071": "\\theta",
    "\uf06c": "\\lambda",
    "\uf077": "\\omega",
    "\uf0b1": "\\pm",
    "\uf0b4": "\\times",
    "\uf0b8": "\\div",
    "\uf0a3": "\\leq",
    "\uf03e": "\\geq",
    "\uf0b9": "\\neq",
    "\uf0d0": "\\infty",
    "\uf0d1": "\\partial",
    "\uf0d2": "\\int",
}

def contains_equation(shape):
  xml = shape._element.xml
  return "m:oMath" in xml or "m:oMathPara" in xml

def is_math_heavy(text: str, threshold: int = 1) -> bool:
  unicode_math = sum(1 for c in text if c in MATH_SYMBOLS)
  symbol_font_hits = sum(1 for c in text if c in SYMBOL_FONT_MATH)
  pattern_hits = len(re.findall(r"\^|\uf053|\uf050|\uf044", text))
  return (unicode_math + symbol_font_hits + pattern_hits) >= threshold

def _get_transform():
  global _xslt_transform
  if _xslt_transform is None:
    if not os.path.exists(XSLT_PATH):
      raise FileNotFoundError(f"OMML2MML.XSL not found at {XSLT_PATH}")
    xslt = etree.parse(XSLT_PATH)
    _xslt_transform = etree.XSLT(xslt)
  return _xslt_transform
  
def _mathml_elem_to_latex_str(elem) -> str:
  tag = elem.tag.split("}")[-1] if "}" in elem.tag else elem.tag
  text = (elem.text or "").strip()
  children = [_mathml_elem_to_latex_str(c) for c in elem]

  if tag == "mfrac":
    return f"\\frac{{{children[0]}}}{{{children[1]}}}" if len(children) == 2 else text
  if tag == "msup":
    return f"{children[0]}^{{{children[1]}}}" if len(children) == 2 else text
  if tag == "msub":
    return f"{children[0]}_{{{children[1]}}}" if len(children) == 2 else text
  if tag == "msqrt":
    return f"\\sqrt{{{' '.join(children)}}}"
  if tag == "mroot":
    return f"\\sqrt[{children[1]}]{{{children[0]}}}" if len(children) == 2 else text
  if tag in ("mrow", "math", "mstyle"):
    return " ".join(children)
  if tag == "msubsup":
    return f"{children[0]}_{{{children[1]}}}^{{{children[2]}}}" if len(children) == 3 else text
  if tag == "mover":
    return f"\\overset{{{children[1]}}}{{{children[0]}}}" if len(children) == 2 else text
  if tag == "munder":
    return f"\\underset{{{children[1]}}}{{{children[0]}}}" if len(children) == 2 else text
  if tag == "msum":
    return f"\\sum"
  if tag in ("mi", "mn", "mo", "mtext"):
    return text
    
  # fallback: just join children or return text
  return " ".join(children) if children else text

def mathml_to_latex(mathml_str: str) -> str | None:
  try:
    root = etree.fromstring(mathml_str.encode())
    return _mathml_elem_to_latex_str(root).strip() or None
  except Exception:
    return re.sub(r"<[^>]+>", "", mathml_str).strip() or None
  
# convert single oMath elem to Latex, OMML -> MathML -> Latex
def _omml_elem_to_latex(math_elem) -> str | None:
  try:
    transform = _get_transform()
  except FileNotFoundError:
    return "".join(math_elem.itertext()).strip() or None
  try:
    omml_str = etree.tostring(math_elem, encoding="unicode")
    omml_tree = etree.fromstring(omml_str.encode())
    mathml_tree = transform(omml_tree)
    mathml_str = etree.tostring(mathml_tree, encoding="unicode")
    return mathml_to_latex(mathml_str)
  except Exception:
    return "".join(math_elem.itertext()).strip() or None

# source is: omml_xslt (full structural conversion succeeded),
# omml_raw (xslt failed, raw text from omml), or none (no math elements)
def omml_to_latex_from_paragraph(paragraph) -> tuple[str | None, str]:
  elem = paragraph._p
  math_elems = elem.findall(f".//{{{OMML_NS}}}oMath")
  if not math_elems:
    return None, "none"
  
  latex_parts = []
  source = "omml_xslt"

  for math_elem in math_elems:
    result = _omml_elem_to_latex(math_elem)
    if result:
      latex_parts.append(result)
    else:
      raw = "".join(math_elem.itertext()).strip()
      if raw:
        latex_parts.append(raw)
        source = "omml_raw"

  if not latex_parts:
    return None, "none"
  return " ".join(latex_parts), source
  
SUPERSCRIPT_MAP = str.maketrans("⁰¹²³⁴⁵⁶⁷⁸⁹", "0123456789")
# wrap nums after ^ in { } for valid latex
def normalize_unicode_math(text: str) -> str:
  for char, latex in SYMBOL_FONT_MAP.items():
    text = text.replace(char, latex)
  text = text.translate(SUPERSCRIPT_MAP)
  # handle ^digits
  text = re.sub(r"\^(\d+)", r"^{\1}", text)
  # handle ^(expression)
  text = re.sub(r"\^(\([^)]+\))", r"^{\1}", text)
  return text

#------------------------------------------------

#------------------------------------------------
# TEXTUAL METRICS (per element)
#------------------------------------------------

# ratio of U+FFFD repl chars to total chars. 0.0 = clean, 1.0 = fully garbled
def detect_garbled_text(text: str) -> float:
  bad_chars = sum(1 for c in text if ord(c) == 65533)
  return bad_chars / max(len(text), 1)

# checks for spacing patterns impl extraction artifacts. 
# 3+ consecutive spaces, punctuation preceded by whitespace, single chars in spaces
def detect_suspicious_spacing(text: str) -> float:
  issues = 0
  if re.search(r" {3,}", text):
    issues += 1
  if re.search(r"\s[.,!?;:]", text):
    issues += 1
  if re.search(r"(?<!\A)\b[a-zA-Z]\b(?!\Z)", text):
    issues += 1
  
  return min(issues / 3, 1.0)

# hard signal: ends on connector char with no closing token
# soft signal: short text that isn't known short form (single word, num, label)
def detect_truncated_text(text: str) -> float:
  stripped = text.strip()
  if len(stripped) < 2:
    return 1.0
  
  hard_truncation = (":", "-", "/", "(", ",", "and", "or", "the", "a")
  if stripped.endswith(tuple(hard_truncation)):
    return 1.0
  
  if (len(stripped.split()) > 8 and stripped[-1].isalpha()):
    return 0.5
  return 0.0

# weighs garbled chars, sus spacing, and truncation
# text quality score for one string, 0.0 (unusable) to 1.0 (clean)
def estimate_text_confidence(text: str) -> float:
  if not text or not text.strip():
    return 0.0
  score = 1.0
  score -= detect_garbled_text(text) * 0.6
  score -= detect_suspicious_spacing(text) * 0.25
  score -= detect_truncated_text(text) * 0.35
  return round(max(0.0, min(score, 1.0)), 4)

#------------------------------------------------
# STRUCTURAL METRICS (per slide)
#------------------------------------------------

# measures how much sort_slide_elements had to reorder extraction
# high score implies already reading order
# called before and after sort_slide_elements
def estimate_reading_order_confidence(elements_before_sort: list, elements_after_sort: list) -> float:
  ids_before = [e["element_id"] for e in elements_before_sort]
  ids_after = [e["element_id"] for e in elements_after_sort]
  if not ids_before:
    return 1.0
  
  matches = sum(a == b for a, b in zip(ids_before, ids_after))
  return round(matches / len(ids_before), 4)

# ratio of element pairs' bboxes substantially overlap
# prevents flagging same bboxes from same parent shape
# 0.0 (no overlaps) to 1.0 (complete overlap)
def detect_overlapping_bboxes(elements: list) -> float:
  def bbox_tuple(b):
    return (b["x"], b["y"], b["width"], b["height"])
  
  def overlaps(a, b):
    if bbox_tuple(a) == bbox_tuple(b):
      return False
    ax2 = a["x"] + a["width"]
    ay2 = a["y"] + a["height"]
    bx2 = b["x"] + b["width"]
    by2 = b["y"] + b["height"]

    # intersection area
    ix = max(0, min(ax2, bx2) - max(a["x"], b["x"]))
    iy = max(0, min(ay2, by2) - max(a["y"], b["y"]))
    intersection = ix * iy

    # only flag if overlap is >20% of smaller element's area
    area_a = a["width"] * a["height"]
    area_b = b["width"] * b["height"]
    smaller = min(area_a, area_b)
    return smaller > 0 and (intersection / smaller) > 0.20
  
  bboxed = [e for e in elements if e.get("bbox")]
  if len(bboxed) < 2:
    return 0.0
  
  total_pairs = 0
  overlapping_pairs = 0

  for i in range(len(bboxed)):
    for j in range(i + 1, len(bboxed)):
      total_pairs += 1
      if overlaps(bboxed[i]["bbox"], bboxed[j]["bbox"]):
        overlapping_pairs += 1
  return round(overlapping_pairs / max(total_pairs, 1), 4)

# count of each element type on the slide
def compute_type_distribution(elements: list) -> dict:
  return dict(Counter(e["type"] for e in elements))

# true if no element has any extractable text content (ex: image only, blank)
def is_slide_empty(elements: list) -> bool:
  for e in elements:
    content = e.get("content")
    if isinstance(content, str) and content.strip():
      return False
    if isinstance(content, list):
      # check table cells for text
      for row in content:
        if any(cell.strip() for cell in row if isinstance(cell, str)):
          return False
  return True

# count of non speakernote elements
def compute_element_density(elements: list) -> int:
  return sum(1 for e in elements if e.get("type") != "speaker_note")

#------------------------------------------------
# METRIC AGGREGATION
#------------------------------------------------

# quality check for every element dict
def compute_element_quality(element: dict) -> dict:
  e_type = element.get("type")
  content = element.get("content")

  if e_type == "table":
    row_count = len(content) if content else 0
    col_count = len(content[0]) if content and content[0] else 0
    empty_cells = sum(
      1 for row in (content or [])
      for cell in row
      if not (isinstance(cell, str) and cell.strip())
    )
    total_cells = row_count * col_count
    return {
      "row_count": row_count,
      "col_count": col_count,
      "empty_cell_ratio": round(empty_cells / max(total_cells, 1), 4),
      "likely_has_header_row": _guess_header_row(content)
    }
  if isinstance(content, str):
    return {
      "char_count": len(content),
      "text_confidence": estimate_text_confidence(content)
    }
  return {}

# first row is likely header if its cells are shorter on avg
# or if all cells are nonempty. soft signal
def _guess_header_row(table_content: list) -> bool:
  if not table_content or len(table_content) < 2:
    return False
  
  first_row = [c for c in table_content[0] if isinstance(c, str) and c.strip()]
  rest = [
    c for row in table_content[1:]
    for c in row
    if isinstance(c, str) and c.strip()
  ]
  if not first_row or not rest:
    return False
  avg_first = mean(len(c) for c in first_row)
  avg_rest = mean(len(c) for c in rest)
  return avg_first < avg_rest

# quality for every slide
def compute_slide_quality(elements_before_sort: list, elements_after_sort: list) -> dict:
  text_scores = []
  for e in elements_after_sort:
    content = e.get("content")
    if isinstance(content, str) and content.strip():
      text_scores.append(estimate_text_confidence(content))

  return {
    "is_empty": is_slide_empty(elements_after_sort),
    "element_density": compute_element_density(elements_after_sort),
    "type_distribution": compute_type_distribution(elements_after_sort),
    "text_confidence": round(mean(text_scores), 4) if text_scores else None,
    "reading_order_confidence": estimate_reading_order_confidence(elements_before_sort, elements_after_sort),
    "overlapping_bbox_ratio": detect_overlapping_bboxes(elements_after_sort)
  }
#------------------------------------------------

##################################################
# DOCX EXTRACTION
##################################################

def omml_to_llm(omml):
  return "will convert later"

def has_page_break_before(paragraph) -> bool:
    """Returns True if this paragraph starts a new page."""
    p = paragraph._p

    # 1. Explicit <w:br w:type="page"/> inside any run
    ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    for br in p.findall(f".//{{{ns}}}br"):
        if br.get(f"{{{ns}}}type") == "page":
            return True

    # 2. <w:pageBreakBefore/> in paragraph properties
    pPr = p.pPr
    if pPr is not None:
        pbr = pPr.find(f"{{{ns}}}pageBreakBefore")
        if pbr is not None:
            # w:val defaults to true if absent; explicit "false"/"0" means off
            val = pbr.get(f"{{{ns}}}val", "true")
            if val.lower() not in ("false", "0"):
                return True

    return False

def has_section_break_before(paragraph) -> bool:
    """Returns True if the PREVIOUS paragraph ended with a section break,
    which forces a new page in most section types."""
    # Section breaks live in <w:pPr><w:sectPr> of the paragraph that
    # PRECEDES the new section — so we check the previous sibling's pPr.
    p = paragraph._p
    prev = p.getprevious()
    if prev is None:
        return False
    ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    pPr = prev.find(f"{{{ns}}}pPr")
    if pPr is None:
        return False
    sectPr = pPr.find(f"{{{ns}}}sectPr")
    if sectPr is None:
        return False
    # Continuous sections don't break to a new page
    pgType = sectPr.find(f"{{{ns}}}type")
    if pgType is not None:
        val = pgType.get(f"{{{ns}}}val", "")
        if val == "continuous":
            return False
    return True

def iter_block_items_with_breaks(doc):
    """
    Yields (block, is_page_break) tuples.
    is_page_break is True when this block starts a new page.
    """
    ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
    body = doc.element.body

    for child in body.iterchildren():
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag

        if tag == "p":
            para = Paragraph(child, doc)

            # Check for explicit page break inside runs
            is_page_break = False
            for br in child.findall(f".//{{{ns}}}br"):
                if br.get(f"{{{ns}}}type") == "page":
                    is_page_break = True
                    break

            # Check for pageBreakBefore in pPr
            if not is_page_break:
                pPr = child.find(f"{{{ns}}}pPr")
                if pPr is not None:
                    pbr = pPr.find(f"{{{ns}}}pageBreakBefore")
                    if pbr is not None:
                        val = pbr.get(f"{{{ns}}}val", "true")
                        if val.lower() not in ("false", "0"):
                            is_page_break = True

            # Check for sectPr inside pPr — paragraph ends a section/page
            # The NEXT paragraph starts a new page, so we yield the break
            # with the next paragraph, handled by carrying a flag
            yield para, is_page_break, tag

        elif tag == "tbl":
            yield Table(child, doc), False, tag

        elif tag == "sectPr":
            # Top-level sectPr directly in body = section/page boundary
            # Signal that the next block starts a new page
            yield None, True, "sectPr"

def get_list_type(paragraph):
  style = paragraph.style.name.lower() if paragraph.style else ""
  
  # numbered lists
  if (
    "list number" in style or
    "number" in style
  ):
    return "numbered"
  
  # bullets
  if (
    "list bullet" in style or
    "bullet" in style
  ):
    return "bullet"
  
  xml = paragraph._element.xml.lower()
  
  if "w:numpr" in xml:
    return "bullet"
  return None

def get_bullet_level_docx(paragraph) -> int:
    """Read indent level from <w:ilvl w:val="N"/> in paragraph XML."""
    try:
        pPr = paragraph._p.pPr
        if pPr is not None:
            ilvl = pPr.find(".//{http://schemas.openxmlformats.org/wordprocessingml/2006/main}ilvl")
            if ilvl is not None:
                val = ilvl.get("{http://schemas.openxmlformats.org/wordprocessingml/2006/main}val")
                if val is not None:
                    return int(val)
        # fallback via python-docx numPr
        try:
            return paragraph._p.pPr.numPr.ilvl.val
        except Exception:
            return 0
    except Exception:
        return 0

def contains_equation(paragraph):
  xml = paragraph._element.xml
  return "m:oMath" in xml or "m:oMathPara" in xml

def extract_equation_xml(paragraph):
  return paragraph._element.xml

def classify_paragraph(paragraph):
    text = clean_text(paragraph.text)

    if not text and not contains_equation(paragraph):
        return None

    style = paragraph.style.name.lower() if paragraph.style else ""
    list_type = get_list_type(paragraph)

    # Check equation first — takes priority over bullet/text classification
    if contains_equation(paragraph):
        latex, eq_source = omml_to_latex_from_paragraph(paragraph)
        if latex:
            # Prepend any plain text in the paragraph before the equation
            plain_prefix = text.strip()
            full_content = f"{plain_prefix} {latex}".strip() if plain_prefix else latex
        else:
            # OMML present but conversion failed — normalize whatever unicode math is there
            full_content = normalize_unicode_math(text) if text else ""
            eq_source = "text_heuristic"
            latex = full_content

        return {
            "type": "equation",
            "content": full_content,
            "equation_latex": latex,
            "equation_source": eq_source
        }

    # Check for hand-typed unicode math (no OMML) in non-bullet text
    if list_type is None and is_math_heavy(text):
        normalized = normalize_unicode_math(text)
        return {
            "type": "equation",
            "content": normalized,
            "equation_latex": normalized,
            "equation_source": "text_heuristic"
        }

    if list_type == "bullet":
        p_type = "bullet"
    elif list_type == "numbered":
        p_type = "numbered"
    elif "title" in style:
        p_type = "title"
    elif "heading" in style:
        p_type = "heading"
    elif "footer" in style:
        p_type = "footer"
    else:
        p_type = "text"

    item = {
        "type": p_type,
        "content": text
    }

    if p_type in ("bullet", "numbered"):
        item["level"] = get_bullet_level_docx(paragraph)

    return item
  

def extract_table(table):
  table_data = []
  for row in table.rows:
    row_data = []

    for cell in row.cells:
      row_data.append(clean_text(cell.text))
    
    table_data.append(row_data)

  return {
      "type": "table",
      "content": table_data
  }

def iter_block_items(doc):
  body = doc.element.body

  for child in body.iterchildren():
    if child.tag.endswith('p'):
      yield Paragraph(child, doc)
    elif child.tag.endswith('tbl'):
      yield Table(child, doc)

def extract_docx(file_path):
  try:
    doc = Document(file_path)

    plain_text = []
    element_counter = 0
    previous_type = None
    number_counter = 1

    pages = []  # list of page element lists
    current_page = []   # elements on current page
    page_index = 0
    next_is_new_page = False

    for block, is_page_break, tag in iter_block_items_with_breaks(doc):
        # top-level sectPr — no content, just signals next block is new page
        if tag == "sectPr":
            if current_page:
                pages.append(current_page)
                current_page = []
                page_index += 1
            continue

        # explicit break on this paragraph
        if is_page_break and current_page:
            pages.append(current_page)
            current_page = []
            page_index += 1

        item = None
        if isinstance(block, Paragraph):
          item = classify_paragraph(block)
        elif isinstance(block, Table):
          item = extract_table(block)

        if item:
          ordered_item = {
            "element_id": f"d_e{element_counter}",
            "type": item["type"],
            "content": item.get("content"),
            "order": element_counter,
            "bbox": None,
          }

          if item["type"] in ("bullet", "numbered") and "level" in item:
            ordered_item["level"] = item["level"]

          ordered_item["quality"] = compute_element_quality(item)

          # carry over any other fields (equation_latex, equation_source, etc.)
          for k, v in item.items():
            if k not in ordered_item:
              ordered_item[k] = v

          current_page.append(ordered_item)

          # collect plain text
          if item["type"] == "table":
            for row in item["content"]:
              cleaned_cells = [cell for cell in row if cell]
              if cleaned_cells:
                plain_text.append("  ".join(cleaned_cells))
          else:
            content = item["content"]
            if item["type"] == "bullet":
              content = f"• {content}"
            elif item["type"] == "numbered":
              if previous_type != "numbered":
                number_counter = 1
              content = f"{number_counter}. {item['content']}"
              number_counter += 1
            if content:
              plain_text.append(content)

          element_counter += 1
          previous_type = item["type"]

    # don't forget the last page
    if current_page:
        pages.append(current_page)

    # build pages output with per-page quality
    pages_output = []
    for pid, page_elements in enumerate(pages):
        text_scores = [
            estimate_text_confidence(e["content"])
            for e in page_elements
            if isinstance(e.get("content"), str) and e["content"].strip()
        ]
        pages_output.append({
            "page_id": pid,
            "elements": page_elements,
            "quality": {
                "is_empty": is_slide_empty(page_elements),
                "element_density": compute_element_density(page_elements),
                "type_distribution": compute_type_distribution(page_elements),
                "text_confidence": round(mean(text_scores), 4) if text_scores else None,
                "reading_order_confidence": 1.0,
                "overlapping_bbox_ratio": 0.0
            }
        })

    os.makedirs("docx", exist_ok=True)
    with open("docx/docx_plain_text.txt", "w", encoding="utf-8") as f:
      f.write("\n".join(plain_text))

    return {
      "file": file_path,
      "pages": pages_output
    }

  except Exception as e:
    import traceback
    traceback.print_exc()
    return None
    
#------------------------------------------------
# DOCX JSON NORMALIZATION
#------------------------------------------------

def normalize_docx(docx_result: dict) -> dict:
  normalized_pages = []

  for page in docx_result.get("pages", []):
    normalized_elements = []

    for elem in page.get("elements", []):
      e_type = elem.get("type")
      content = elem.get("content")

      # drop footers — not curriculum content
      if e_type == "footer":
        continue

      if e_type == "table":
        normalized_elements.append({
          "type": "table",
          "content": content,
          "quality": elem.get("quality")
        })
        continue

      if not isinstance(content, str) or not content.strip():
        continue

      # normalize title → heading to match pptx/pdf
      if e_type == ("title", "subtitle"):
        e_type = "heading"

      norm = {
        "type": e_type,
        "content": content
      }

      if e_type == "bullet" and elem.get("level") is not None:
        norm["level"] = elem["level"]

      if e_type == "numbered":
        norm["type"] = "bullet"
        norm["numbered"] = True
        if elem.get("level") is not None:
          norm["level"] = elem["level"]

      if elem.get("ocr_source") is not None:
        norm["ocr_source"] = elem["ocr_source"]
        if elem.get("ocr_confidence") is not None:
          norm["ocr_confidence"] = elem["ocr_confidence"]

      if e_type == "equation":
        if elem.get("equation_latex"):
          norm["equation_latex"] = elem["equation_latex"]
        if elem.get("equation_source"):
          norm["equation_source"] = elem["equation_source"]

      if elem.get("quality"):
        norm["quality"] = elem["quality"]

      normalized_elements.append(norm)

    normalized_pages.append({
      "page_id": page["page_id"],
      "elements": normalized_elements,
      "quality": page.get("quality")
    })

  return {
    "file": docx_result.get("file"),
    "pages": normalized_pages
  }

##################################################
# POWERPOINT EXTRACTION
##################################################
# file_path = "/content/Sample Presentation.pptx"

def get_bullet_type(paragraph):
  pPr = paragraph._p.pPr
  if pPr is not None:
  
    if pPr.find(".//{*}buAutoNum") is not None:
      return "numbered"
    
    if pPr.find(".//{*}buChar") is not None:
      return "bullet"
    
    if pPr.find(".//{*}buNone") is not None:
      return None
  
  if paragraph.level > 0:
    return "bullet"
  
  # bullet types in PPT XML
  return None

def sort_slide_elements(elements):
  title_bottom = 0

  # First pass: find title bottom (this must be its own loop)
  for elem in elements:
    if elem["type"] == "title" and elem["bbox"]:
        bbox = elem["bbox"]
        if "y" in bbox:
            title_bottom = max(title_bottom, bbox["y"] + bbox["height"])
        else:
            title_bottom = max(title_bottom, bbox["y1"])

  def _get_center_x(elem):
    bbox = elem.get("bbox", {})
    if "x" in bbox:
        return bbox["x"] + bbox["width"] / 2
    elif "x0" in bbox:
        return (bbox["x0"] + bbox["x1"]) / 2
    return None

  body_centers = [
      _get_center_x(elem)
      for elem in elements
      if elem["type"] != "title" and elem.get("bbox") and _get_center_x(elem) is not None
  ]

    # Detect column split point by finding the largest gap between sorted x-positions
  col_split = None
  if len(body_centers) >= 2:
    sorted_xs = sorted(set(body_centers))
    if len(sorted_xs) >= 2:
      gaps = [
          (sorted_xs[i+1] - sorted_xs[i], sorted_xs[i], sorted_xs[i+1])
          for i in range(len(sorted_xs) - 1)
      ]
      largest_gap = max(gaps, key=lambda g: g[0])
      gap_size, gap_left, gap_right = largest_gap

      SLIDE_WIDTH = 9144000  # standard EMU width
      if gap_size > SLIDE_WIDTH * 0.10:
        col_split = (gap_left + gap_right) / 2

  def reading_order_key(elem):
    bbox = elem.get("bbox")

    if not bbox:
      return (999, 999, 999999999)

    # Support both bbox schemas: x/y/width/height and x0/y0/x1/y1
    if "x" in bbox:
        x = bbox["x"]
        y = bbox["y"]
        w = bbox["width"]
    else:
        x = bbox["x0"]
        y = bbox["y0"]
        w = bbox["x1"] - bbox["x0"]

    cx = x + w / 2

    if elem["type"] == "title":
      return (0, 0, y)

    # Assign column based on dynamic split point
    if col_split is not None:
      column = 0 if x < col_split else 1
    else:
      column = 0

    return (1, column, y)

  return sorted(elements, key=reading_order_key)

def is_inside_bbox(inner, outer):
  return(
    inner["x"] >= outer["x"] and
    inner["y"] >= outer["y"] and
    inner["x"] + inner["width"] <= outer["x"] + outer["width"] and
    inner["y"] + inner["height"] <= outer["y"] + outer["height"]
  )

def detect_shape_type(shape, paragraph=None):
  if paragraph is not None:

    # check if paragraph contains math xml, priority over bullet type
    para_xml = paragraph._p.xml
    if "m:oMath" in para_xml or "m:oMathPara" in para_xml:
      return "equation"
    
    bullet_type = get_bullet_type(paragraph)
    if bullet_type:
      # check if it contains heavy math unicode
      text = paragraph.text or ""
      if is_math_heavy(text):
        return "equation"
      return bullet_type
    
  try:
    if shape.is_placeholder:
      placeholder = shape.placeholder_format.type

      if placeholder == PP_PLACEHOLDER.TITLE:
        return "title"
      
      if placeholder == PP_PLACEHOLDER.SUBTITLE:
        return "subtitle"
  except:
    pass
    
  if contains_equation(shape):
    return "equation"
  
  return "text"

def classify_shape(shape, slide_num, element_counter, plain_text, metric_text, slide=None):
  elements = []
  bbox = get_bbox(shape)

  # detect slides with images and alt text/descriptions (if any)
  if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
    # Attempt OCR on the image
    try:
      ocr_elements = extract_image_elements_pptx(shape, slide_num, element_counter)
    except Exception as e:
      print(f"  [OCR] pptx image OCR failed: {e}")
      ocr_elements = []

    if ocr_elements:
      # OCR succeeded — attach alt text/description as metadata and return
      for elem in ocr_elements:
        elem["element_id"] = f"s{slide_num}_e{element_counter}"
        elem["order"] = element_counter
        element_counter += 1
      return ocr_elements, element_counter

    # OCR returned nothing — fall back to metadata-only image element
    elements.append({
      "element_id": f"s{slide_num}_e{element_counter}",
      "type": "image",
      "content": None,
      "caption": None,
      "caption_source": None,
      "order": element_counter,
      "bbox": bbox,
      "quality": {
        "has_caption": False,
        "caption_source": None
      }
    })
    element_counter += 1
    return elements, element_counter

  # tables
  if shape.has_table:
    table_data = []
    tbl = shape.table
    row_count = len(tbl.rows)
    col_count = len(tbl.columns)

    for r in range(row_count):
      row_data = []

      for c in range(col_count):
        cell = tbl.cell(r, c)
        cell_text = clean_text(cell.text)

        if not cell_text and cell.text_frame:
          run_texts = []

          for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
              txt = clean_text(run.text)
              if txt:
                run_texts.append(txt)
          cell_text = " ".join(run_texts)

        row_data.append(cell_text)

      table_data.append(row_data)

      # plain text
      cleaned_cells = [c for c in row_data if c]
      if cleaned_cells:
        row_text = "  ".join(cleaned_cells)
        plain_text.append(row_text)
        metric_text.append(row_text)

    elements.append({
        "element_id": f"s{slide_num}_e{element_counter}",
        "type": "table",
        "content": table_data,
        "order": element_counter,
        "bbox": bbox
    })
    element_counter += 1
    return elements, element_counter

  if shape.has_text_frame:

    number_counter = 1

    for paragraph in shape.text_frame.paragraphs:
      text = clean_text(paragraph.text)

      if not text:
        continue

      p_type = detect_shape_type(shape, paragraph)

      element = {
        "element_id": f"s{slide_num}_e{element_counter}",
        "type": p_type,
        "content": text,
        "order": element_counter,
        "bbox": bbox
      }

      if p_type == "bullet":
        element["level"] = paragraph.level
        serialized_text = f"• {text}"

      elif p_type == "numbered":
        serialized_text = f"{number_counter}.  {text}"
        number_counter += 1

      elif p_type == "equation":
        latex, eq_source = omml_to_latex_from_paragraph(paragraph)
        if latex:
          # prepend plain text in paragraph before equation
          plain_prefix = clean_text(paragraph.text).strip()
          full_content = f"{plain_prefix} {latex}".strip() if plain_prefix else latex
          element["content"] = full_content
          element["equation_latex"] = latex
          element["equation_source"] = eq_source
        else:
          # hand typed unicode math (no omml), store as is
          normalized = normalize_unicode_math(text)
          element["content"] = normalized
          element["equation_latex"] = normalized
          element["equation_source"] = "text_heuristic"
        serialized_text = element["content"]

      else:
        serialized_text = text

      elements.append(element)
      element_counter += 1

      # store plain text
      plain_text.append(serialized_text)
      metric_text.append(text)

  return elements, element_counter

def iter_shapes(shapes):
  for shape in shapes:
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
      yield from iter_shapes(shape.shapes)
    else:
      yield shape

def extract_speaker_notes(slide, slide_num, element_counter, plain_text, metric_text):
  notes = []

  if not slide.has_notes_slide:
    return notes, element_counter

  for shape in slide.notes_slide.shapes:
    # skip slide image placeholder; has .text but not a note
    if not shape.has_text_frame:
      continue

    try:
      if shape.is_placeholder:
        ph_type = shape.placeholder_format.type
        if ph_type != PP_PLACEHOLDER.BODY:
          continue
    except:
      continue

    text = clean_text(shape.text_frame.text)

    if not text:
      continue
      
    notes.append({
        "element_id": f"s{slide_num}_e{element_counter}",
        "type": "speaker_note",
        "content": text,
        "order": element_counter,
        "bbox": None
      })
    element_counter += 1

    plain_text.append(text)
    metric_text.append(text)

  return notes, element_counter

def extract_pptx(file_path):
  try:
    presentation = Presentation(file_path)

    result = {"file": file_path, "slides": []}

    plain_text = []
    # capture only text for cer/wer metrics
    metric_text = []

    for slide_num, slide in enumerate(presentation.slides, start=1):
      slide_data = {"slide_id": slide_num, "elements": []}

      element_counter = 0

      for shape in iter_shapes(slide.shapes):
        elements, element_counter = classify_shape(shape, slide_num, element_counter, plain_text, metric_text, slide)

        if elements:
          slide_data["elements"].extend(elements)

      notes, element_counter = extract_speaker_notes(slide, slide_num, element_counter, plain_text, metric_text)
      slide_data["elements"].extend(notes)

      # capture pre sort order
      elements_before_sort = list(slide_data["elements"])

      slide_data["elements"] = sort_slide_elements(slide_data["elements"])

      # attach per element quality
      for e in slide_data["elements"]:
        if e.get("type") == "image_text":
          continue
        e["quality"] = compute_element_quality(e)

      slide_data["quality"] = compute_slide_quality(elements_before_sort, slide_data["elements"])

      result["slides"].append(slide_data)

    # save plain text
    #os.makedirs("pptx", exist_ok=True)
    #with open("pptx/pptx_plain_text.txt", "w", encoding="utf-8") as f:
      #f.write("\n".join(plain_text))

    #with open("pptx/pptx_metric_text.txt", "w", encoding="utf-8") as f:
      #f.write("\n".join(metric_text))

    return result

  except Exception as e:
    print(f"Error: {e}")
    return {"error": str(e)}
  
#------------------------------------------------
# PPTX JSON NORMALIZATION
#------------------------------------------------

def normalize_pptx(pptx_result: dict) -> dict:
  normalized_pages = []

  for slide in pptx_result.get("slides", []):
    normalized_elements = []

    for elem in slide.get("elements", []):
      e_type = elem.get("type")
      content = elem.get("content")

      # omit speaker notes — not curriculum content
      #if e_type == "speaker_note":
        #continue

      # omit images with no extractable text
      #if e_type == "image":
        #continue

      if e_type == "table":
        normalized_elements.append({
          "type": "table",
          "content": content,
          "quality": elem.get("quality")
        })
        continue

      if not isinstance(content, str) or not content.strip():
        continue

      # normalize title/subtitle → heading
      if e_type in ("title", "subtitle"):
        e_type = "heading"

      norm = {
        "type": e_type,
        "content": content
      }

      if e_type == "bullet" and elem.get("level") is not None:
        norm["level"] = elem["level"]

      # numbered → bullet with numbered flag
      if e_type == "numbered":
        norm["type"] = "bullet"
        norm["numbered"] = True

      if e_type == "equation":
        if elem.get("equation_latex"):
          norm["equation_latex"] = elem["equation_latex"]
        if elem.get("equation_source"):
          norm["equation_source"] = elem["equation_source"]

      if elem.get("ocr_source") is not None:
        norm["ocr_source"] = elem["ocr_source"]
        if elem.get("ocr_confidence") is not None:
          norm["ocr_confidence"] = elem["ocr_confidence"]

      if elem.get("quality"):
        norm["quality"] = elem["quality"]

      normalized_elements.append(norm)

    normalized_pages.append({
      "page_id": slide["slide_id"],
      "elements": normalized_elements,
      "quality": slide.get("quality")
    })

  return {
    "file": pptx_result.get("file"),
    "pages": normalized_pages
  }

##################################################
# PDF EXTRACTION
##################################################

# file_path = "CSE330 Midterm Review PDF.pdf"

def ocr_pipeline(file_path):
  print("This is not a text PDF!")
  return [], []

def overlaps(word, bbox):
  return ()

def get_bullet_level(text, x0=None, page_width=None):
    stripped = text.strip()

    char_level = None
    if stripped.startswith("■"):
        char_level = 2
    elif stripped.startswith("○"):
        char_level = 1
    elif stripped.startswith(("●", "-", "•")):
        char_level = 0

    indent_level = None
    if x0 is not None and page_width is not None:
        indent_ratio = x0 / page_width
        if indent_ratio > 0.20:
            indent_level = 2
        elif indent_ratio > 0.12:
            indent_level = 1
        elif indent_ratio > 0.05:
            indent_level = 0

    if char_level is not None and indent_level is not None:
        return char_level
    if char_level is not None:
        return char_level
    if indent_level is not None and stripped.startswith(("-", "–", "*")):
        return indent_level

    return None

# table arrays are out of order
def split_into_cols(region_words, page_width=None):
  if not region_words:
    return[]

  valid_words = [
      w for w in region_words
      if isinstance(w, dict) and "x0" in w
  ]

  if not valid_words:
    return []

  x0s = np.array([w["x0"] for w in valid_words])

  # detect spread
  spread = np.std(x0s)

  if page_width is None or spread < page_width * 0.25:
    return [valid_words]

  if page_width is not None:
    split_x = page_width / 2
  else:
    split_x = np.percentile(x0s, 50)

  left = [w for w in valid_words if w["x0"] < split_x]
  right = [w for w in valid_words if w["x0"] >= split_x]

  # avoid fake column detection
  if (
      len(left) < 0.2 * len(valid_words)
      or len(right) < 0.2 * len(valid_words)
      or len(left) < 3
      or len(right) < 3
  ):
    return [valid_words]

  return [left, right]

def group_lines(words, tol=3):
  lines = []

  for w in sorted(words, key=lambda w: w["top"]):
    placed = False

    for line in lines:
      if abs(line["top"] - w["top"]) < tol:
        line["words"].append(w)
        placed = True
        break

    if not placed:
      lines.append({"top": w["top"], "words": [w]})

  return lines

def is_real_table(t, page_width, page_height):
    bbox = t["bbox"]
    data = t["data"]

    flat = [cell for row in data for cell in row]
    if not flat:
        return False

    empty_ratio = sum(1 for c in flat if str(c).strip() == "") / len(flat)

    bbox_w = bbox[2] - bbox[0]
    bbox_h = bbox[3] - bbox[1]
    is_full_page = bbox_w > page_width * 0.85 and bbox_h > page_height * 0.85

    # reject full-page mostly-empty grids
    if is_full_page and empty_ratio > 0.85:
        return False

    # reject mostly-empty grids of any size
    if empty_ratio > 0.90:
        return False

    # reject 2-column "tables" — almost certainly layout boxes, not data
    num_cols = max(len(row) for row in data) if data else 0
    if num_cols <= 2:
        return False

    # reject if non-empty cells contain long prose (layout box, not data cell)
    non_empty = [str(c).strip() for c in flat if str(c).strip()]
    if non_empty:
        avg_len = sum(len(c) for c in non_empty) / len(non_empty)
        if avg_len > 60:
            return False

    return True


 # BORROWING PPTX METRICS
def pdf_bbox_to_metrics_fmt(elem):
    """Convert pdf-style bbox to the format detect_overlapping_bboxes expects."""
    b = elem.get("bbox")
    if not b:
        return None
    return {
        **elem,
        "bbox": {
            "x": b["x0"],
            "y": b["y0"],
            "width": b["x1"] - b["x0"],
            "height": b["y1"] - b["y0"]
        }
    }

def extract_pdf(file_path):
  try:
    plain_text_pdf = []
    text_pages = 0
    pages_output = []

    with pdfplumber.open(file_path) as pdf:
      # for /content/CSE330 Midterm Review.pdf
      start_page = 0
      end_page = len(pdf.pages)-1

      # debugging
      print("Total pages:", len(pdf.pages))
      print("Extracting pages:", start_page, "to", end_page)

      # collect tables
      table_groups = []

      #-------------------------
      # Table extraction
      #-------------------------

      pages = pdf.pages[start_page: end_page + 1]

      for page_index, page in enumerate(pages, start=start_page):
        tables = page.find_tables()

        # find and extract table data first so that text extraction
        # does not occur twice over the table
        for table in tables:
          extracted = table.extract()
          if not extracted:
            continue

          normalized = [
              [cell if cell is not None else "" for cell in row]
              for row in extracted
          ]

          table_groups.append({
              "page": page_index,
              "bbox": table.bbox,
              "data": normalized
          })

      table_groups.sort(key=lambda t: (t["page"], t["bbox"][1]))

      # merge tables that span across page breaks
      def merge_split_tables(groups):
        merged = []

        for t in groups:
          data = t["data"]

          if not data:
            # or not isinstance(data, list):
            continue

          if not merged:
            merged.append(t)
            continue

          last = merged[-1]

          if not last["data"]:
            merged.append(t)
            continue

          # column mismatch means new table
          if len(data[0]) != len(last["data"][0]):
            merged.append(t)
            continue

          overlap = sum(
              1 for a, b in zip(data[0], last["data"][0]) if a == b
          )
          similarity = overlap / max(len(data[0]), 1)

          # check top of next page
          is_next_page = t["page"] == last["page"] + 1
          near_top = t["bbox"][1] < 120

          # detect headerless table continuation
          if similarity > 0.7 or (is_next_page and near_top):
            last["data"].extend(data)
          else:
            merged.append(t)

        return merged

      merged_tables = merge_split_tables(table_groups)

      # filter out false tables detected from decorative borders/lines
      merged_tables = [
          t for t in merged_tables
          if is_real_table(t, pdf.pages[t["page"]].width, pdf.pages[t["page"]].height)
      ]

      #-------------------------
      # Main page loop
      # process text sections and positioning,
      # filtered to exclude table sections
      #-------------------------
      all_page_images = convert_from_path(file_path, dpi=150)

      pages = pdf.pages[start_page:end_page + 1]
      for page_index, page in enumerate(pages, start=start_page):
        # extract words
        words = page.extract_words()
        if words:
          text_pages += 1

        # use original tables for filtering, not merged
        original_page_tables = [
            t for t in table_groups
            if t["page"] == page_index
        ]

        # uses raw unvalidated detections, wrong
        #table_bboxes = [t["bbox"] for t in original_page_tables]

        # only exclude words from tables that actually passed is_real_table
        validated_page_tables = [
          t for t in merged_tables
          if t["page"] == page_index
        ]

        table_bboxes = [t["bbox"] for t in validated_page_tables]

        # filter words inside tables
        # disabled, fix later
        #filtered_words = words
        def word_in_table(w, table_bboxes):
          for bbox in table_bboxes:
            if w["x0"] >= bbox[0] and w["top"] >= bbox[1] and w["x1"] <= bbox[2] and w["bottom"] <= bbox[3]:
              return True
          return False

        filtered_words = [w for w in words if not word_in_table(w, table_bboxes)]

        if not filtered_words:
          pages_output.append({
              "page_id": page_index,
              "elements": [],
              "quality": {
                  "is_empty": True,
                  "element_density": 0,
                  "type_distribution": {},
                  "text_confidence": None,
                  "reading_order_confidence": 1.0,
                  "overlapping_bbox_ratio": 0.0
              }
          })
          continue

        filtered_words.sort(key=lambda w: w["top"])

        #-------------------------
        # Region splitting
        # text in columns must be grouped correctly, and
        # there may be one or multiple column sections
        # throughout the file
        #-------------------------

        regions = []
        current = []
        last_top = None

        #filtered_words = sorted(filtered_words, key = lambda w: (w["top"], w["x0"]))

        for w in filtered_words:
          if last_top is None:
            current = [w]
          elif abs(w["top"] - last_top) > 15:
            regions.append(current)
            current = [w]
          else:
            current.append(w)
          last_top = w["top"]

        if current:
          regions.append(current)

        def is_noise_region(region):
          if len(region) >= 2:
              return False
          # single-word region: keep if it's in the left 75% of the page (likely real text)
          # discard if it's far right (likely a floating image label)
          return region[0]["x0"] > page.width * 0.75

        regions = [r for r in regions if not is_noise_region(r)]

        #-------------------------
        # Line extraction
        #-------------------------
        page_lines = [] # store structured lines

        for region_idx, region in enumerate(regions):

          columns = split_into_cols(region, page_width=page.width)
          if not columns:
            continue

          columns = [col for col in columns
                     if isinstance(col, list) and len(col) > 0]

          if not columns:
            continue

          columns.sort(
              key=lambda col: np.median([w["x0"] for w in col if isinstance(w, dict) and "x0" in w])
          )
          #columns.sort(key = lambda col: np.median([w["x0"] for w in col if isinstance(w, dict)]))

          # read regions in correct order
          for col_idx, col_words in enumerate(columns):
            grouped = group_lines(col_words, tol = 3)

            for line in sorted(grouped, key = lambda l: l["top"]):
              line_words = sorted(line["words"], key = lambda x: x["x0"])

              #-------------------------
              # Text reconstruction
              #-------------------------
              #line_text = " ".join(w["text"] for w in line_words).strip()

              gaps = [
                  line_words[i]["x0"] - line_words[i-1]["x1"]
                  for i in range(1, len(line_words))
              ]
              gaps = [g for g in gaps if g > 0]

              if gaps:
                median_gap = np.median(gaps)
              else:
                median_gap = 2

              threshold = max(1.5, median_gap * 0.6)

              line_text = ""
              prev_x1 = None

              for w in line_words:
                if prev_x1 is None:
                  line_text += w["text"]
                else:
                  gap = w["x0"] - prev_x1

                  if gap > threshold:
                    line_text += " " + w["text"]
                  else:
                    line_text += w["text"]
                prev_x1 = w["x1"]

              page_lines.append({
                  "text": line_text,
                  "top": line["top"],
                  "x0": min(w["x0"] for w in line_words),
                  "x1": max(w["x1"] for w in line_words),
                  "bottom": max(w["bottom"] for w in line_words),
                  "region": region_idx,
                  "column": col_idx,
                  "is_bullet": get_bullet_level(line_text, x0=min(w["x0"] for w in line_words), page_width=page.width) is not None,
                  "bullet_level": get_bullet_level(line_text, x0=min(w["x0"] for w in line_words), page_width=page.width)              })

        #-------------------------
        # Paragraph merging
        #-------------------------
        def should_merge(prev, curr, page_width):
          vertical_gap = curr["top"] - prev["top"]
          proximity = vertical_gap < 22  # enough for typical line spacing

          prev_level = prev.get("bullet_level")
          curr_level = curr.get("bullet_level")
          prev_text  = prev["text"].rstrip()
          curr_text  = curr["text"].lstrip()

          if not proximity:
              return False

          # never merge two bullet lines — new bullet char always starts fresh
          if prev_level is not None and curr_level is not None:
              return False

          # never absorb a continuation into something that isn't a bullet
          if prev_level is None and curr_level is not None:
              return False

          # ── continuation line into a bullet ──────────────────────────────────────
          if prev_level is not None and curr_level is None:
              # prev ended mid-phrase (no sentence-ending punctuation)
              no_terminal = not prev_text.endswith((".", "!", "?", ":"))
              if not no_terminal:
                  return False

              # the continuation x0 must sit to the right of the bullet marker
              # (the marker is always the leftmost part of the bullet line)
              # and must not reach the next indent threshold, which would make
              # it a sub-bullet rather than a wrapped line
              x0_prev   = prev["x0"]
              x0_curr   = curr["x0"]
              x0_ratio  = x0_curr / page_width

              is_to_right_of_marker = x0_curr > x0_prev

              # make sure it isn't sitting at a known bullet-start indent
              # (if it were a bullet it would have a char — but be safe)
              not_a_new_indent_level = not any([
                  x0_ratio > 0.20,   # would be level 2
              ])

              # the line continues: starts lowercase/digit, or prev ended softly
              prev_ends_soft  = prev_text.endswith((",", ";", "—", "-"))
              starts_lower    = bool(curr_text) and (
                  curr_text[0].islower() or curr_text[0].isdigit()
              )

              return is_to_right_of_marker and not_a_new_indent_level and (
                  starts_lower or prev_ends_soft
              )

          # ── two plain text lines ──────────────────────────────────────────────────
          indent = abs(curr["x0"] - prev["x0"]) < 10
          return indent
    

        paragraphs = []
        current_para = []
        page_width = page.width

        #for line in sorted(page_lines, key = lambda x: x["top"]):
        for line in sorted(page_lines, key=lambda x: (x["region"], x["column"], x["top"])):

          if not current_para:
            current_para.append(line)
            continue

          if should_merge(current_para[-1], line, page_width):
            current_para.append(line)
          else:
            paragraphs.append(current_para)
            current_para = [line]

        if current_para:
          paragraphs.append(current_para)

        #-------------------------
        # classify element type
        #-------------------------
        def classify_type(para, page_height):
          first = para[0]
          text = " ".join(l["text"] for l in para).strip()
          is_near_top = first["top"] < page_height * 0.15
          is_short = len(text) < 80
          bullet_level = first["bullet_level"]

          if is_math_heavy(text):
            return "equation", 0

          if bullet_level is not None:
              return "bullet", bullet_level
          if is_near_top and is_short:
              return "heading", 0
          return "text", 0


        #-------------------------
        # Text stream
        #-------------------------
        page_stream = []
        seen_texts = set()
        order = 0

        page_tables = [t for t in merged_tables if t["page"] == page_index]
        inserted = set()

        # OCR 
        page_img_pil = all_page_images[page_index - start_page]
        image_elements = extract_image_elements(page, page_img_pil, page_index, dpi=150)

        # collect all elements with y0 for sorting
        raw_elements = []

        for para in paragraphs:
            text = " ".join(l["text"] for l in para).strip()
            if not text or text in seen_texts:
                continue
            seen_texts.add(text)
            elem_type, elem_level = classify_type(para, page.height)

            elem = {
                "element_id": f"p{page_index}_e{order}",
                "type": elem_type,
                "content": text,
                "order": order,
                "bbox": {
                    "x0": min(l["x0"] for l in para),
                    "y0": para[0]["top"],
                    "x1": max(l["x1"] for l in para),
                    "y1": para[-1]["bottom"]
                }
            }

            if elem_type == "bullet" and elem_level is not None:
                elem["level"] = elem_level
            
            if elem_type == "equation":
                elem["equation_latex"] = normalize_unicode_math(text)
                elem["equation_source"] = "text_heuristic"

            elem["quality"] = compute_element_quality({
                  "type": elem_type,
                  "content": text
                })

            raw_elements.append({"_y0": para[0]["top"], **elem})

        for t in page_tables:
            bbox = t["bbox"]
            raw_elements.append({
                "element_id": f"p{page_index}_e_tbl",  # placeholder; overwritten below
                "type": "table",
                "content": t["data"],
                "order": -1,   # also placeholder
                "_y0": bbox[1],
                "bbox": {"x0": bbox[0], "y0": bbox[1], "x1": bbox[2], "y1": bbox[3]},
                "quality": compute_element_quality({"type": "table", "content": t["data"]})
            })

        # capture order before sort for reading order confidence
        """elements_before_sort = [
            {"element_id": e.get("element_id", ""), **e}
            for e in raw_elements
        ]"""

        for i, e in enumerate(raw_elements):
          e["_presort_idx"] = i

        ids_before_sort = list(range(len(raw_elements)))

        # sort text and tables first, images later
        raw_elements.sort(key=lambda e: e["_y0"])

        # after sort, presort indices tell us new order
        ids_after_sort_presort = [e["_presort_idx"] for e in raw_elements]

        # reading order confidence: how many are still in original relative order
        matches = sum(a == b for a, b in zip(ids_before_sort, ids_after_sort_presort))
        reading_order_confidence = round(matches / max(len(ids_before_sort), 1), 4)

        # append image elements at the end, inserted by y0 but after any
        # overlapping text element
        for img_elem in image_elements:
            img_y0 = img_elem["bbox"]["y0"]

            # insert after the last text element that starts before the image's top edge
            insert_at = 0
            for i, e in enumerate(raw_elements):
                if e["_y0"] < img_y0:
                    insert_at = i + 1

            raw_elements.insert(insert_at, {"_y0": img_y0, **img_elem})

        for elem in raw_elements:
            elem.pop("_y0")
            elem.pop("_presort_idx", None)  # clean up temp field
            elem["element_id"] = f"p{page_index}_e{order}"
            elem["order"] = order
            page_stream.append(elem)

            if elem["type"] == "table":
                plain_text_pdf.append(
                    "\n".join("\t".join(str(c) for c in row) for row in elem["content"])
                )
            elif elem.get("ocr_source"):
                plain_text_pdf.append(f"[OCR] {elem['content']}")
            else:
                plain_text_pdf.append(elem["content"])

            order += 1

        # adapt page_stream bboxes for overlap detection
        adapted_stream = [pdf_bbox_to_metrics_fmt(e) for e in page_stream if e.get("bbox")]

        text_scores = [
            estimate_text_confidence(e["content"])
            for e in page_stream
            if isinstance(e.get("content"), str) and e["content"].strip()
            and e.get("type") not in ("image",)
        ]

        # ── slide-level quality ──────────────────────────────────────────────────────
        elements_after_sort = page_stream

        slide_quality = {
            "is_empty": is_slide_empty(page_stream),
            "element_density": compute_element_density(page_stream),
            "type_distribution": compute_type_distribution(page_stream),
            "text_confidence": round(mean(text_scores), 4) if text_scores else None,
            "reading_order_confidence": reading_order_confidence,
            "overlapping_bbox_ratio": detect_overlapping_bboxes(adapted_stream)
        }

        pages_output.append({
            "page_id": page_index,
            "elements": page_stream,
            "quality": slide_quality
        })

    if text_pages == 0:
      # fully scanned PDF
      for page_index, page in enumerate(pdf.pages[start_page:end_page + 1], start=start_page):
        page_img_pil = all_page_images[page_index - start_page]
        image_elements = extract_image_elements(page, page_img_pil, page_index, dpi=150)
        # assign element_id and order
        for i, elem in enumerate(image_elements):
          elem["element_id"] = f"p{page_index}_e{i}"
          elem["order"] = i
        pages_output.append({
          "page_id": page_index,
          "elements": image_elements,
          "quality": {
            "is_empty": len(image_elements) == 0,
            "element_density": len(image_elements),
            "type_distribution": {"image_text": len(image_elements)} if image_elements else {},
            "text_confidence": round(float(np.mean([e["ocr_confidence"] for e in image_elements])), 4) if image_elements else None,
            "reading_order_confidence": 1.0,
            "overlapping_bbox_ratio": 0.0
          }
        })
      result = {"file": file_path, "pages": pages_output}
      return result, plain_text_pdf
    
    result = {
      "file": file_path,
      "pages": pages_output
    }

    return result, plain_text_pdf

  except Exception as e:
    import traceback
    traceback.print_exc()
    return f"Error extracting PDF text: {e}", []
  
#------------------------------------------------
# PDF JSON NORMALIZATION
# strip element_id, order, bbox
#------------------------------------------------
def normalize_pdf(pdf_result: dict) -> dict:
  normalized_pages = []

  for page in pdf_result.get("pages", []):
    normalized_elements = []

    for elem in page.get("elements", []):
      e_type = elem.get("type")
      content = elem.get("content")

      if e_type == "table":
        normalized_elements.append({
          "type": "table",
          "content": content,
          "quality": elem.get("quality")
        })
        continue

      if not isinstance(content, str) or not content.strip():
        continue

      norm = {
        "type": e_type,
        "content": content
      }

      if e_type == "bullet" and elem.get("level") is not None:
        norm["level"] = elem["level"]

      if e_type == "equation":
        if elem.get("equation_latex"):
          norm["equation_latex"] = elem["equation_latex"]
        if elem.get("equation_source"):
          norm["equation_source"] = elem["equation_source"]

      if elem.get("ocr_source") is not None:
        norm["ocr_source"] = elem["ocr_source"]
        if elem.get("ocr_confidence") is not None:
          norm["ocr_confidence"] = elem["ocr_confidence"]

      if elem.get("quality"):
        norm["quality"] = elem["quality"]

      normalized_elements.append(norm)

    normalized_pages.append({
      "page_id": page["page_id"],
      "elements": normalized_elements,
      "quality": page.get("quality")
    })

  return {
    "file": pdf_result.get("file"),
    "pages": normalized_pages
  }

##################################
# MAIN
##################################

from lesson_paths import lesson_path                      # ← add

def run_text_extraction(file_path: str, lesson_id: str) -> dict: 
    """
    Runs extraction + normalization for one upload, writes
    normalized_output.json into that lesson's folder, and returns it
    so the caller (server.py) can pass it straight into the next stage
    without a redundant disk read.
    """
    file_type = Path(file_path).suffix.lower()

    if file_type == ".pdf":
        pdf_result, plain_text_pdf = extract_pdf(file_path)
        normalized = normalize_pdf(pdf_result)
    elif file_type == ".docx":
        docx_result = extract_docx(file_path)
        normalized = normalize_docx(docx_result)
    elif file_type == ".pptx":
        pptx_result = extract_pptx(file_path)
        normalized = normalize_pptx(pptx_result)
    else:
        raise ValueError(f"Unsupported file type: {file_type}")

    out_path = lesson_path(lesson_id, "normalized_output.json", create_dir=True)
    with open(out_path, "w", encoding="utf-8") as f:
        json.dump(normalized, f, indent=2, ensure_ascii=False)

    return normalized


if __name__ == "__main__":                        
    file_path = input("Please enter file path: ").strip()
    lesson_id = input("Lesson id (blank = 'manual-test'): ").strip() or "manual-test"
    run_text_extraction(file_path, lesson_id)
    print(f"Done. Output written to backend/lessons/{lesson_id}/normalized_output.json")


"""
file_path = input("Please enter file path: ").strip()
print("File path is:", file_path)

file_type = Path(file_path).suffix.lower()
print("File type is:", file_type)

# file_path = "/content/advanced sample document.pdf"
# /content/CSE330 Midterm Review.pdf
if file_type == '.pdf':
  pdf_result, plain_text_pdf = extract_pdf(file_path)
  os.makedirs("pdf", exist_ok=True)
  with open("pdf/pdf_output.json", "w", encoding="utf-8") as f:
    json.dump(pdf_result, f, indent=2, ensure_ascii=False)

  with open("pdf/pdf_output.txt", "w", encoding="utf-8") as f:
    f.write("\n\n".join(plain_text_pdf))

  normalized_pdf_result = normalize_pdf(pdf_result)
  with open("normalized_output.json", "w", encoding="utf-8") as f:
    json.dump(normalized_pdf_result, f, indent=2, ensure_ascii=False)

  print(json.dumps(pdf_result, indent=2, ensure_ascii=False))

# /content/Sample docx document.docx
elif file_type == '.docx':
  docx_result = extract_docx(file_path)
  os.makedirs("docx", exist_ok=True)
  with open("docx/docx_output.json", "w", encoding="utf-8") as f:
    json.dump(docx_result, f, indent=2, ensure_ascii=False)

  normalized_docx_result = normalize_docx(docx_result)
  with open("normalized_output.json", "w", encoding="utf-8") as f:
    json.dump(normalized_docx_result, f, indent=2, ensure_ascii=False)

  print(json.dumps(docx_result, indent=2, ensure_ascii=False))

# /content/Sample Presentation.pptx
# /content/CSE330 Midterm Review.pptx
elif file_type == '.pptx':
  pptx_result = extract_pptx(file_path)
  os.makedirs("pptx", exist_ok=True)
  with open("pptx/pptx_output.json", "w", encoding="utf-8") as f:
    json.dump(pptx_result, f, indent=2, ensure_ascii=False)
  print(json.dumps(pptx_result, indent=2, ensure_ascii=False))

  normalized_pptx_result = normalize_pptx(pptx_result)
  with open("normalized_output.json", "w", encoding="utf-8") as f:
    json.dump(normalized_pptx_result, f, indent=2, ensure_ascii=False)

else:
  print(f"Unsupported file type: {file_type}")
"""